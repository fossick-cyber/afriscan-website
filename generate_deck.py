"""Generate AfriScan sales presentation deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# Brand colors
DARK = RGBColor(0x0F, 0x11, 0x17)
ACCENT = RGBColor(0x4F, 0x8F, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED = RGBColor(0xFF, 0x45, 0x3A)
ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
GREEN = RGBColor(0x34, 0xC7, 0x59)
SUBTLE = RGBColor(0x8B, 0x8F, 0xA3)

OUTPUT = Path("/opt/favhousecheck/website/assets/AfriScan-Presentation.pptx")


def set_slide_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=14, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_accent_bar(slide, top):
    shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(top), Inches(1.5), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    add_text(slide, 1, 1.5, 11, 1.2,
             "AfriScan", size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 2.7, 11, 0.6,
             "by Afridrone", size=20, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_text(slide, 1.5, 3.8, 10, 1.2,
             "AI-Powered Pipeline Corridor Scanning\n& Encroachment Detection for Africa",
             size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 2, 5.5, 9, 0.5,
             "Protect your infrastructure. Know what's on the ground.",
             size=16, color=SUBTLE, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # SLIDE 2: The Problem
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "The Problem", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    problems = [
        ("Community Encroachment",
         "New buildings, settlements, and farming creep into pipeline corridors and concession areas. Without monitoring, encroachment goes undetected until it becomes a crisis."),
        ("Costly Disputes",
         "Undetected land use changes lead to compensation claims, legal battles, and project delays. A single relocation dispute can cost millions."),
        ("Safety & Compliance",
         "People living near pipelines and powerlines face real danger. Regulators hold operators accountable for encroachment in their corridor."),
    ]
    for i, (title, desc) in enumerate(problems):
        y = 1.8 + i * 1.7
        add_text(slide, 1, y, 10, 0.5, title, size=20, bold=True, color=ACCENT)
        add_text(slide, 1, y + 0.45, 10, 0.8, desc, size=14, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════
    # SLIDE 3: The Solution
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "The Solution: AfriScan", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    add_text(slide, 1, 1.8, 11, 1.0,
             "Draw a pipeline route. Click scan.\nGet a complete encroachment report — in minutes.",
             size=22, color=WHITE, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Draw Route", "Draw on map, upload KML,\nor select existing route"),
        ("2", "AI Scans", "3 AI models analyse\nsatellite imagery"),
        ("3", "Risk Analysis", "Every 500m segment\nclassified by risk"),
        ("4", "Report", "PDF report, Google Earth,\nGIS exports"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = 1 + i * 3
        # Circle with number
        shape = slide.shapes.add_shape(
            9, Inches(x + 0.7), Inches(3.3), Inches(0.7), Inches(0.7))  # oval
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text(slide, x, 4.2, 2.5, 0.5, title, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, x, 4.7, 2.5, 0.8, desc, size=12, color=SUBTLE, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # SLIDE 4: AI Detection
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "Multi-Model AI Detection", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    models = [
        ("Google Open Buildings", "398 buildings", "Pre-mapped building footprints\nfrom Google's global dataset", GREEN),
        ("RAMP Deep Learning", "+537 additional", "Semantic segmentation model trained\nfor African rural structures", ACCENT),
        ("Custom YOLO", "Fine-tuned", "Train on your own imagery\nfor maximum accuracy", ORANGE),
    ]
    for i, (name, stat, desc, color) in enumerate(models):
        x = 0.8 + i * 4
        add_text(slide, x, 1.8, 3.5, 0.5, name, size=20, bold=True, color=color)
        add_text(slide, x, 2.3, 3.5, 0.5, stat, size=28, bold=True, color=WHITE)
        add_text(slide, x, 2.9, 3.5, 0.8, desc, size=13, color=LIGHT_GRAY)

    add_text(slide, 1, 4.2, 11, 0.8,
             "Combined: 935 buildings detected on a 4.6km test corridor",
             size=18, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, 1, 5.0, 11, 1.0,
             "Purpose-built for Africa — detects mud huts, thatched roofs,\nzinc compounds, and informal structures that generic models miss.",
             size=14, color=SUBTLE, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # SLIDE 5: Risk Analysis
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "Risk Segment Analysis", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    add_text(slide, 1, 1.6, 11, 0.6,
             "Pipeline divided into 500m segments, each classified by encroachment density",
             size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Risk table header
    risk_data = [
        ("0 – 500m", "HIGH", "16 buildings", RED),
        ("500 – 1000m", "MEDIUM", "5 buildings", ORANGE),
        ("1000 – 1500m", "HIGH", "34 buildings", RED),
        ("1500 – 2000m", "HIGH", "41 buildings", RED),
        ("2000 – 2500m", "MEDIUM", "2 buildings", ORANGE),
        ("2500 – 3000m", "HIGH", "8 buildings", RED),
        ("3000 – 3500m", "HIGH", "18 buildings", RED),
    ]
    y_start = 2.4
    add_text(slide, 1.5, y_start - 0.4, 3, 0.3, "Segment", size=11, bold=True, color=SUBTLE)
    add_text(slide, 4.5, y_start - 0.4, 2, 0.3, "Risk", size=11, bold=True, color=SUBTLE)
    add_text(slide, 7, y_start - 0.4, 3, 0.3, "Buildings", size=11, bold=True, color=SUBTLE)

    for i, (seg, risk, count, color) in enumerate(risk_data):
        y = y_start + i * 0.45
        add_text(slide, 1.5, y, 3, 0.4, seg, size=13, color=WHITE)
        add_text(slide, 4.5, y, 2, 0.4, risk, size=13, bold=True, color=color)
        add_text(slide, 7, y, 3, 0.4, count, size=13, color=LIGHT_GRAY)

    add_text(slide, 1, 5.8, 11, 0.6,
             "Real data from Pande test corridor, Mozambique — 4.6km pipeline route",
             size=12, color=SUBTLE, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # SLIDE 6: Deliverables
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "What You Get", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    deliverables = [
        ("Professional PDF Report",
         "Executive summary, risk rating, building photos with GPS,\ncoordinate tables, risk segment analysis"),
        ("Interactive Web Map",
         "Colour-coded risk segments, toggleable building layers,\nsatellite imagery, click any detection for details"),
        ("Google Earth Export (KMZ)",
         "Open in Google Earth — buildings as placemarks,\nrisk segments as colour-coded lines, full route overlay"),
        ("GIS Data (GeoPackage / GeoJSON)",
         "Every building with distance to pipeline, buffer zone,\nconfidence score, source model — ready for your GIS team"),
        ("Historical Change Tracking",
         "Compare scans over time — new buildings highlighted,\ntrend analysis for progressive encroachment"),
    ]
    for i, (title, desc) in enumerate(deliverables):
        y = 1.6 + i * 1.05
        add_text(slide, 1, y, 5, 0.4, title, size=16, bold=True, color=ACCENT)
        add_text(slide, 1, y + 0.35, 10, 0.6, desc, size=12, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════
    # SLIDE 7: Industries
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "Industries We Serve", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    industries = [
        ("Oil & Gas", "Pipeline corridor monitoring,\nright-of-way encroachment,\nconstruction zone surveys"),
        ("Mining", "Concession boundary monitoring,\nartisanal mining detection,\nsettlement tracking"),
        ("Power & Utilities", "Transmission line corridors,\nvegetation encroachment,\nexclusion zone monitoring"),
        ("Telecoms", "Fibre route surveys,\ncell tower site monitoring,\ncoverage planning"),
    ]
    for i, (name, desc) in enumerate(industries):
        x = 0.5 + i * 3.2
        add_text(slide, x, 1.8, 2.8, 0.5, name, size=22, bold=True, color=ACCENT)
        add_text(slide, x, 2.5, 2.8, 1.5, desc, size=13, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════
    # SLIDE 8: Why AfriScan
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.8, "Why AfriScan", size=36, bold=True, color=WHITE)
    add_accent_bar(slide, 1.2)

    points = [
        "Africa-trained AI — detects structures generic models miss",
        "One-click scanning — no GIS expertise required",
        "3 AI models working together for maximum coverage",
        "Risk-rated pipeline segments — prioritise where to act",
        "Professional reports ready for regulators and stakeholders",
        "Google Earth export — results your field teams can use immediately",
        "Historical tracking — monitor encroachment trends over time",
        "Mobile-ready — run scans and review results from any device",
    ]
    add_bullet_list(slide, 1.2, 1.6, 10, 4.5, points, size=16, color=WHITE)

    # ═══════════════════════════════════════════════════════
    # SLIDE 9: Contact
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text(slide, 1, 1.5, 11, 1.0,
             "AfriScan", size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 2.5, 11, 0.6,
             "by Afridrone", size=22, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_text(slide, 1.5, 3.5, 10, 1.0,
             "Protect your pipeline corridors and concessions\nwith AI-powered encroachment monitoring.",
             size=20, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, 1, 5.0, 11, 0.5,
             "afri-scan.com", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 5.5, 11, 0.5,
             "afridr.one", size=16, color=SUBTLE, align=PP_ALIGN.CENTER)

    # Save
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT} ({OUTPUT.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    build()
